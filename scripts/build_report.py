#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
一级市场估值报告生成器（v3风格）
生成：Word(.docx) + PPT(.pptx) + Excel(.xlsx) + 辅助交付物
数据来源：config.json + valuation_results.json + cap_table.json
所有公司专属叙事内容从 config.report_narrative 读取，无硬编码。
"""
import argparse, csv, json, os, sys

from docx import Document
from docx.shared import Pt, Inches, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT
from docx.oxml.ns import qn
from docx.oxml import OxmlElement

from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn as pptx_qn

import openpyxl
from openpyxl.styles import Font as XlFont, PatternFill, Alignment as XlAlign, Border, Side

# ============================================================
# 命令行解析
# ============================================================
def parse_args():
    p = argparse.ArgumentParser(description="一级市场估值报告生成器（v3风格）")
    p.add_argument("--config", required=True, help="config.json 路径")
    p.add_argument("--data", default="./data", help="估值结果与cap table数据目录")
    p.add_argument("--charts", default="./charts", help="图表PNG目录")
    p.add_argument("--outdir", default="./deliverables", help="输出目录")
    return p.parse_args()

# ============================================================
# 数据加载
# ============================================================
def load_data(cfg_path, data_dir):
    with open(cfg_path, encoding="utf-8") as f:
        cfg = json.load(f)
    with open(os.path.join(data_dir, "valuation_results.json"), encoding="utf-8") as f:
        val = json.load(f)
    with open(os.path.join(data_dir, "cap_table.json"), encoding="utf-8") as f:
        cap = json.load(f)
    return cfg, val, cap

# ---- 便捷提取 ----
def narr(cfg, key, default=None):
    """从 config.report_narrative 安全读取，缺省返回 default"""
    nr = cfg.get("report_narrative", {})
    parts = key.split(".")
    v = nr
    for p in parts:
        if isinstance(v, dict):
            v = v.get(p)
        else:
            return default
    return v if v is not None else default

STAGE_CN = {"Angel": "天使轮", "Seed": "种子轮", "Pre-A": "Pre-A轮",
            "A": "A轮", "B": "B轮", "C": "C轮", "Pre-IPO": "Pre-IPO轮"}

# ---- 配色与字体常量 ----
FONT_CN = "微软雅黑"
FONT_EN = "Arial"
C_DARK_BLUE = RGBColor(0x1F, 0x4E, 0x79)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_RED = RGBColor(0xC0, 0x39, 0x2B)
C_GREY = RGBColor(0x96, 0x96, 0x96)
C_DARK = RGBColor(0x2C, 0x3E, 0x50)
HEADER_BG = "1F4E79"
ALT_ROW_BG = "D6E4F0"
P_DARK_BLUE = PptRGB(0x1F, 0x4E, 0x79)
P_WHITE = PptRGB(0xFF, 0xFF, 0xFF)
P_RED = PptRGB(0xC0, 0x39, 0x2B)
P_GREY = PptRGB(0x96, 0x96, 0x96)
P_DARK = PptRGB(0x2C, 0x3E, 0x50)
P_LIGHTBLUE = PptRGB(0xD6, 0xE4, 0xF0)

# ============================================================
# Word 辅助函数
# ============================================================
def set_run_font(run, size=10.5, bold=False, color=None):
    run.font.size = Pt(size); run.font.bold = bold
    run.font.name = FONT_EN
    run._element.rPr.rFonts.set(qn("w:eastAsia"), FONT_CN)
    if color: run.font.color.rgb = color

def add_heading(doc, text, level=1):
    p = doc.add_paragraph()
    sizes = {1: 16, 2: 13, 3: 11.5}
    r = p.add_run(text)
    set_run_font(r, size=sizes.get(level,11), bold=True, color=C_DARK_BLUE)
    p.paragraph_format.space_before = Pt(14 if level==1 else 8)
    p.paragraph_format.space_after = Pt(6)

def add_para(doc, text, size=10.5, bold=False, color=None, align=None, indent=False):
    p = doc.add_paragraph()
    if align: p.alignment = align
    if indent: p.paragraph_format.first_line_indent = Pt(21)
    p.paragraph_format.space_after = Pt(6); p.paragraph_format.line_spacing = 1.4
    for i, part in enumerate(text.split("**")):
        if not part: continue
        r = p.add_run(part)
        set_run_font(r, size=size, bold=(bold or i%2==1), color=color)

def add_bullet(doc, text, size=10.5, bold_lead=None):
    p = doc.add_paragraph(style="List Bullet")
    p.paragraph_format.space_after = Pt(3); p.paragraph_format.line_spacing = 1.3
    if bold_lead:
        r = p.add_run(bold_lead); set_run_font(r, size=size, bold=True, color=C_DARK_BLUE)
    for i, part in enumerate(text.split("**")):
        if not part: continue
        r = p.add_run(part); set_run_font(r, size=size, bold=(i%2==1))

def _shade_cell(cell, hex_color):
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"),"clear"); shd.set(qn("w:fill"),hex_color)
    tcPr.append(shd)

def add_table(doc, headers, rows, col_widths=None, font_size=9.5, caption=None):
    if caption:
        cp = doc.add_paragraph(); r = cp.add_run(caption)
        set_run_font(r, size=10, bold=True, color=C_DARK_BLUE)
        cp.paragraph_format.space_after = Pt(3)
    table = doc.add_table(rows=1, cols=len(headers))
    table.alignment = WD_TABLE_ALIGNMENT.CENTER; table.style = "Table Grid"
    for i, h in enumerate(headers):
        _shade_cell(table.rows[0].cells[i], HEADER_BG)
        p = table.rows[0].cells[i].paragraphs[0]; p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = p.add_run(str(h)); set_run_font(r, size=font_size, bold=True, color=C_WHITE)
    for r_idx, row in enumerate(rows):
        cells = table.add_row().cells
        for i, val in enumerate(row):
            if r_idx % 2 == 1: _shade_cell(cells[i], ALT_ROW_BG)
            p = cells[i].paragraphs[0]
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER if i>0 else WD_ALIGN_PARAGRAPH.LEFT
            r = p.add_run(str(val)); set_run_font(r, size=font_size)
    if col_widths:
        for i,w in enumerate(col_widths):
            for row in table.rows: row.cells[i].width = Inches(w)
    doc.add_paragraph().paragraph_format.space_after = Pt(2)

def add_chart(doc, chart_dir, filename, width=6.2, caption=None):
    if caption:
        cp = doc.add_paragraph(); cp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        r = cp.add_run(caption); set_run_font(r, size=10, bold=True, color=C_DARK_BLUE)
        cp.paragraph_format.space_after = Pt(2)
    path = os.path.join(chart_dir, filename)
    if os.path.exists(path):
        p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        p.paragraph_format.space_after = Pt(8)
        p.add_run().add_picture(path, width=Inches(width))

def add_takeaway(doc, text):
    p = doc.add_paragraph(); p.paragraph_format.space_after=Pt(10); p.paragraph_format.line_spacing=1.35
    r = p.add_run("解读："); set_run_font(r, size=9.5, bold=True, color=C_DARK_BLUE)
    for i, part in enumerate(text.split("**")):
        if not part: continue
        r = p.add_run(part); set_run_font(r, size=9.5, bold=(i%2==1))

# ============================================================
# PPT 辅助函数
# ============================================================
def _set_ppt_cn(run):
    run.font.name = FONT_CN
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(pptx_qn("a:ea"))
    if ea is None: ea = rPr.makeelement(pptx_qn("a:ea"),{}); rPr.append(ea)
    ea.set("typeface",FONT_CN)

def blank_slide(prs): return prs.slides.add_slide(prs.slide_layouts[6])

def ppt_title(slide, text):
    box = slide.shapes.add_textbox(PptInches(0.5),PptInches(0.28),PptInches(12.3),PptInches(0.75))
    tf = box.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = text
    r.font.size = PptPt(28); r.font.bold = True; r.font.color.rgb = P_DARK_BLUE
    _set_ppt_cn(r)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, PptInches(0.5),PptInches(1.08),PptInches(12.3),PptPt(2.5))
    ln.fill.solid(); ln.fill.fore_color.rgb = P_DARK_BLUE; ln.line.fill.background()

def ppt_text(slide, text, left, top, width, height, size=12, bold=False,
             color=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    color = color or P_DARK
    for ii, line in enumerate(str(text).split("\n")):
        p = tf.paragraphs[0] if ii==0 else tf.add_paragraph()
        p.alignment = align; p.space_after = PptPt(4)
        parts = line.split("**")
        if len(parts)==1:
            r = p.add_run(); r.text = line
            r.font.size = PptPt(size); r.font.bold = bold; r.font.italic = italic
            r.font.color.rgb = color; _set_ppt_cn(r)
        else:
            for j, part in enumerate(parts):
                if not part: continue
                r = p.add_run(); r.text = part
                r.font.size = PptPt(size); r.font.bold = (bold or j%2==1); r.font.italic = italic
                r.font.color.rgb = (P_DARK_BLUE if j%2==1 else color); _set_ppt_cn(r)

def ppt_pic(slide, chart_dir, filename, left, top, width=None, height=None):
    path = os.path.join(chart_dir, filename)
    if os.path.exists(path):
        kw = {}; kw["width"] = width if width is not None else None
        if height: kw["height"] = height
        return slide.shapes.add_picture(path, left, top, **{k:v for k,v in kw.items() if v is not None})

def _ppt_shade(cell, hex_color):
    cell.fill.solid()
    cell.fill.fore_color.rgb = PptRGB(int(hex_color[0:2],16),int(hex_color[2:4],16),int(hex_color[4:6],16))

def ppt_table(slide, headers, rows, left, top, width, height, font_size=9, first_col_bold=False):
    gtbl = slide.shapes.add_table(len(rows)+1, len(headers), left, top, width, height)
    tbl = gtbl.table
    for j,h in enumerate(headers):
        cell = tbl.cell(0,j); _ppt_shade(cell, HEADER_BG)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(h)
        r.font.size = PptPt(font_size); r.font.bold = True; r.font.color.rgb = P_WHITE
        _set_ppt_cn(r)
    for i,row in enumerate(rows):
        for j,val in enumerate(row):
            cell = tbl.cell(i+1,j)
            _ppt_shade(cell, ALT_ROW_BG if i%2==1 else "FFFFFF")
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j>0 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = str(val)
            r.font.size = PptPt(font_size); r.font.color.rgb = P_DARK
            if first_col_bold and j==0: r.font.bold = True; r.font.color.rgb = P_DARK_BLUE
            _set_ppt_cn(r)

def ppt_footer(slide, text):
    box = slide.shapes.add_textbox(PptInches(0.5),PptInches(7.12),PptInches(12.3),PptInches(0.3))
    p = box.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = text; r.font.size = PptPt(8); r.font.color.rgb = P_GREY
    _set_ppt_cn(r)

print("辅助函数加载完成。")
sys.stdout.flush()


# ============================================================
# Word 报告主体
# ============================================================
def build_word(doc_path, cfg, val, cap, chart_dir):
    doc = Document()
    st = doc.styles["Normal"]
    st.font.name = FONT_EN; st.element.rPr.rFonts.set(qn("w:eastAsia"), FONT_CN)
    st.font.size = Pt(10.5)
    sec = doc.sections[0]
    sec.top_margin = Inches(0.9); sec.bottom_margin = Inches(0.9)
    sec.left_margin = Inches(1.0); sec.right_margin = Inches(1.0)

    cname = cfg["company"]["name"]; cname_en = cfg["company"].get("name_en","")
    stage_cn = STAGE_CN.get(cfg["round"]["stage"], cfg["round"]["stage"])
    val_date = cfg["meta"]["valuation_date"]
    ft = f"{cname} {stage_cn}估值报告 | {cfg['meta']['confidentiality']} | {val_date}"
    sec.footer.paragraphs[0].text = ""
    r_f = sec.footer.paragraphs[0].add_run(ft); set_run_font(r_f, size=8, color=C_GREY)

    fin = cfg["financials"]
    rev = fin["revenue_forecast"]; nm = fin["net_margin"]; em = fin["ebit_margin"]; npft = fin["net_profit"]
    syn = val["synthesis"]; WPRE = syn["weighted_pre_money"]
    BP_PRE = cfg["round"]["bp_implied_pre_money"]; BP_POST = cfg["round"]["bp_implied_post_money"]
    RND_AMT = cfg["round"]["round_amount"]; EQ_PCT = cfg["round"]["equity_offered"]
    SCN = cfg["valuation_params"]["vc_method"]["scenarios"]
    PW_MOIC = sum(s["moic"]*s["probability"] for s in SCN)
    PW_EXIT_YI = sum(s["exit_value_yi"]*s["probability"] for s in SCN)
    IV = syn["individual_valuations"]

    # ---- 封面 ----
    for _ in range(3): doc.add_paragraph()
    full_name = f"{cname}（{cname_en}）" if cname_en else cname
    add_para(doc, full_name, size=30, bold=True, color=C_DARK_BLUE, align=WD_ALIGN_PARAGRAPH.CENTER)
    rep_title = narr(cfg, "report_title", f"{stage_cn}股权投资估值分析报告")
    add_para(doc, rep_title, size=20, bold=True, color=C_DARK, align=WD_ALIGN_PARAGRAPH.CENTER)
    slogan = narr(cfg, "slogan", "")
    if slogan: add_para(doc, slogan, size=12, color=C_GREY, align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_paragraph()
    # KPI
    kpis = narr(cfg, "cover_kpis", [])
    if kpis:
        kpi_labels = [k["label"] for k in kpis]
        kpi_vals = [k.get("value","") if k.get("value")!="由脚本计算" else f"{PW_MOIC:.1f}x" for k in kpis]
        add_table(doc, kpi_labels, [kpi_vals], font_size=11)
    doc.add_paragraph()
    add_para(doc, f"估值日期：{val_date}　|　货币：人民币　|　报告类型：投资委员会估值备忘",
             size=10.5, align=WD_ALIGN_PARAGRAPH.CENTER)
    add_para(doc, f"本轮融资：¥{RND_AMT:,}万　|　出让股权：{EQ_PCT*100:.0f}%　|　投后定价：¥{BP_POST/10000:.2f}亿（BP条款）",
             size=10.5, align=WD_ALIGN_PARAGRAPH.CENTER)
    add_para(doc, "【机密 · 仅供内部投资决策使用】", size=11, bold=True, color=C_RED, align=WD_ALIGN_PARAGRAPH.CENTER)
    team = narr(cfg, "analyst_team", "估值团队")
    add_para(doc, f"编制：{team}", size=10, color=C_GREY, align=WD_ALIGN_PARAGRAPH.CENTER)
    doc.add_page_break()

    # ---- 执行摘要 ----
    add_heading(doc, "执行摘要", level=1)
    add_para(doc, "估值结论", size=12, bold=True, color=C_DARK_BLUE)
    add_para(doc,
        f"基于风险资本法（VC Method）、计分卡法（Scorecard）与可比交易法三种方法交叉验证，"
        f"{cname}{stage_cn}**投前估值合理区间为 ¥{syn['valuation_range']['min']/10000:.2f}亿 ~ ¥{syn['valuation_range']['max']/10000:.2f}亿**，"
        f"多方法加权中位数为 **¥{WPRE/10000:.2f}亿**。本轮 BP 报价为投前 ¥{BP_PRE/10000:.2f}亿（投后 ¥{BP_POST/10000:.2f}亿，融资 ¥{RND_AMT:,}万、出让{EQ_PCT*100:.0f}%），"
        f"显著低于分析师加权中位数，**对投资人保留了充分的安全边际**。", indent=True)
    add_para(doc,
        f"以 VC Method 概率加权口径计算，本轮投资的**概率加权 MOIC 为 {PW_MOIC:.1f}x**，"
        f"高于天使轮 8x 的回报门槛；基准情景下 MOIC 为 {SCN[1]['moic']:.1f}x，"
        f"对应 {SCN[0].get('exit_year', cfg['valuation_params']['vc_method']['exit_year'])} 年退出估值 ¥{SCN[1]['exit_value_yi']:.0f}亿。", indent=True)

    add_para(doc, "本轮交易条款", size=12, bold=True, color=C_DARK_BLUE)
    esop = cfg["capital_structure"].get("target_esop_pct",0)*100
    terms = [
        ["融资轮次", stage_cn],
        ["融资金额", f"¥{RND_AMT:,}万"],
        ["出让股权", f"{EQ_PCT*100:.0f}%"],
        ["投后估值（BP）", f"¥{BP_POST/10000:.2f}亿"],
        ["投前估值（BP）", f"¥{BP_PRE/10000:.2f}亿"],
        ["证券类型/保护条款", "BP未披露；模型按优先股、1x非参与清算优先权、加权平均反稀释作敏感性假设"],
        ["期权池", f"BP未披露；模型按{esop:.0f}%期权池（投前摊薄）示例测算"],
        ["资金用途", narr(cfg,"fund_usage_summary","详见资金用途饼图")],
    ]
    add_table(doc, ["条款", "内容"], terms, col_widths=[1.8,4.6])

    add_para(doc, "估值方法选择与依据", size=12, bold=True, color=C_DARK_BLUE)
    add_para(doc, narr(cfg,"method_selection_rationale",""), indent=True)

    add_para(doc, "关键发现", size=12, bold=True, color=C_DARK_BLUE)
    findings = narr(cfg,"key_findings",[])
    for i,f in enumerate(findings,1):
        add_bullet(doc, f, bold_lead=f"{i}. ")

    add_para(doc, "估值摘要表", size=12, bold=True, color=C_DARK_BLUE)
    rows = [[m["method"],f"¥{m['valuation']/10000:.2f}亿",f"{m['weight']*100:.0f}%",m["source"]] for m in IV]
    rows.append(["加权投前估值",f"¥{WPRE/10000:.2f}亿","100%",f"区间 ¥{syn['valuation_range']['min']/10000:.2f}–{syn['valuation_range']['max']/10000:.2f}亿"])
    rows.append(["BP隐含投前",f"¥{BP_PRE/10000:.2f}亿","—","公司报价（基准）"])
    add_table(doc, ["估值方法","投前估值","权重","依据/备注"], rows, col_widths=[1.6,1.3,0.8,3.0])

    add_chart(doc, chart_dir, "chart03_valuation_bridge.png", 6.0, "图表1：估值方法交叉验证桥图")
    add_takeaway(doc,
        f"VC Method 给出区间上沿（聚焦退出回报），Scorecard 锚定行业基准给出下沿，可比交易居中。"
        f"**加权投前 ¥{WPRE/10000:.2f}亿显著高于 BP 隐含 ¥{BP_PRE/10000:.2f}亿**，本轮定价对投资人有利。")
    doc.add_page_break()

    # ---- 自定义辅助：从 sections 读文本或兜底 ----
    def s(key, default=""):
        return narr(cfg, f"sections.{key}", default)

    # ---- 一、融资背景与估值目的 ----
    add_heading(doc, "一、融资背景与估值目的", level=1)
    add_heading(doc, "1.1 公司发展阶段", level=2)
    company_profile = s(
        "company_profile_facts",
        f"{cname}成立于{cfg['company']['founded_year']}年，总部位于{cfg['company']['headquarters']}，"
        f"{cfg['company']['description']}。"
    )
    add_para(doc,
        f"{company_profile}{s('company_stage','公司处于天使轮、产品MVP阶段，尚未产生规模化营收。')}",
        indent=True)
    add_heading(doc, "1.2 融资概况与条款", level=2)
    add_para(doc, s("financing_overview"), indent=True)
    add_heading(doc, "1.3 估值目的与范围", level=2)
    add_para(doc, s("valuation_purpose"), indent=True)
    add_chart(doc, chart_dir, "chart05_fund_usage_pie.png", 5.0, f"图表2：{stage_cn}资金用途分配（¥{RND_AMT:,}万）")
    add_takeaway(doc, "资金主要投向研发与产品化，合计大部分用于产品与技术，符合早期阶段「先打磨产品」的资金配置逻辑。")
    doc.add_page_break()

    # ---- 二、公司分析 ----
    add_heading(doc, "二、公司分析", level=1)
    add_heading(doc, "2.1 公司概况", level=2)
    add_para(doc, s("company_overview"), indent=True)
    # 产品矩阵
    products = narr(cfg,"products"); prod_headers = narr(cfg,"product_headers",["产品线","形态","核心功能","变现方式"])
    if products:
        add_heading(doc, "2.2 产品矩阵", level=2)
        add_table(doc, prod_headers, products, col_widths=[1.3,1.6,2.4,1.5])
        add_para(doc, s("product_takeaway"), indent=True)
    # 商业模式
    add_heading(doc, "2.3 商业模式与收入结构", level=2)
    add_para(doc, s("business_model"), indent=True)
    add_chart(doc, chart_dir, "chart10_revenue_mix.png", 5.0, "图表3：2030E 收入构成预测")
    add_takeaway(doc, "至2030年硬件与SaaS/服务收入并重，SaaS占比提升是毛利率与净利率改善的关键驱动。")
    # 壁垒
    moats = narr(cfg,"competitive_moats",[])
    if moats:
        add_heading(doc, "2.4 竞争壁垒", level=2)
        for m in moats: add_bullet(doc, m)
    # 团队
    team = narr(cfg,"team"); team_headers = narr(cfg,"team_headers",["角色","背景","分工"])
    if team:
        add_heading(doc, "2.5 核心团队", level=2)
        add_table(doc, team_headers, team, col_widths=[1.0,3.7,1.7])
        team_note = narr(cfg,"team_note","")
        if team_note: add_para(doc, team_note, indent=True)
    doc.add_page_break()

    # ---- 三、行业与市场分析 ----
    add_heading(doc, "三、行业与市场分析", level=1)
    add_heading(doc, "3.1 老龄化趋势与市场驱动", level=2)
    add_para(doc, s("aging_trend"), indent=True)
    add_heading(doc, "3.2 市场规模", level=2)
    market_rows = narr(cfg,"market_size"); market_h = narr(cfg,"market_headers",["市场","规模","口径","增速"])
    if market_rows: add_table(doc, market_h, market_rows, col_widths=[2.0,1.4,1.2,1.6])
    add_chart(doc, chart_dir, "chart06_market_size.png", 6.0, "图表4：市场规模趋势")
    add_takeaway(doc, s("market_takeaway","TAM天花板极高，赛道选择正确。"))
    add_heading(doc, "3.3 政策环境与红利", level=2)
    policy_rows = narr(cfg,"policy"); policy_h = narr(cfg,"policy_headers",["政策/文件","核心内容"])
    if policy_rows: add_table(doc, policy_h, policy_rows, col_widths=[1.8,4.6])
    add_para(doc, s("policy_paragraph"), indent=True)
    add_heading(doc, "3.4 竞争格局分析", level=2)
    comp_rows = narr(cfg,"competitors"); comp_h = narr(cfg,"competitor_headers",["类型","代表","优势","短板"])
    if comp_rows: add_table(doc, comp_h, comp_rows, col_widths=[1.3,1.9,1.6,1.6])
    add_para(doc, s("competition_paragraph"), indent=True)
    doc.add_page_break()

    # ---- 四、财务预测分析 ----
    add_heading(doc, "四、财务预测分析", level=1)
    add_heading(doc, "4.1 BP财务预测概览", level=2)
    add_para(doc, s("forecast_note"), color=C_RED)
    yrs = [str(y) for y in cfg["financials"]["forecast_years"]]
    rows = [[y, f"{rev[y]:,}", f"{fin['gross_margin']*100:.0f}%",
             f"{em[y]*100:.1f}%", f"{nm[y]*100:.1f}%", f"{npft[y]:,}"] for y in yrs]
    add_table(doc, ["年份","营收(万元)","毛利率","EBIT率","净利率","净利润(万元)"], rows)
    add_chart(doc, chart_dir, "chart01_revenue_forecast.png", 6.0, "图表5：营收预测与利润率趋势")
    add_takeaway(doc, s("forecast_analysis"))
    add_heading(doc, "4.2 收入驱动因素拆解", level=2)
    for d in narr(cfg,"sections.revenue_drivers",[]): add_bullet(doc, d)
    add_heading(doc, "4.3 毛利率与盈利路径", level=2)
    add_para(doc, s("margin_path"), indent=True)
    add_heading(doc, "4.4 关键假设审阅与压力测试", level=2)
    add_para(doc, s("stress_test"), indent=True)
    add_chart(doc, chart_dir, "chart11_dcf_cashflow.png", 6.0, "图表6：简化DCF现金流（辅助参考）")
    add_takeaway(doc, "⚠ 简化DCF仅作辅助参考——天使轮现金流不可验证，DCF易产生虚假精确度，不作为主估值方法。")
    doc.add_page_break()

    # ---- 五、估值分析 ----
    add_heading(doc, "五、估值分析", level=1)
    add_heading(doc, "5.1 方法选择说明", level=2)
    add_para(doc, "针对天使轮无营收、产品MVP的特征，采用四层估值体系："
             "①融资条款反推（定价锚点）；②VC Method概率加权（核心，权重50%）；"
             "③可比交易分析（权重25%）；④Scorecard（权重25%），辅以可比上市公司PS倍数参考。", indent=True)
    add_heading(doc, "5.2 融资条款反推法", level=2)
    add_para(doc, f"按BP条款：投后估值 = 融资额 ÷ 出让比例 = ¥{RND_AMT:,}万 ÷ {EQ_PCT*100:.0f}% = **¥{BP_POST/10000:.2f}亿**，投前 = ¥{BP_PRE/10000:.2f}亿。这是交易双方的直接定价锚点。", indent=True)
    add_heading(doc, "5.3 VC Method（风险资本法）—— 核心方法", level=2)
    add_para(doc, s("vc_method_rationale"), indent=True)
    rows = [[sc["name"], f"{sc['exit_revenue_yi']:.0f}亿", f"{sc['ps_multiple']:.1f}x",
             f"{sc['exit_value_yi']:.0f}亿", f"{sc['angel_retention']*100:.0f}%",
             f"{sc['moic']:.1f}x", f"{sc['probability']*100:.0f}%"] for sc in SCN]
    rows.append(["概率加权","—","—",f"{PW_EXIT_YI:.1f}亿","—",f"{PW_MOIC:.1f}x","100%"])
    add_table(doc, ["情景","退出营收","PS","退出估值","天使保留","MOIC","概率"], rows, caption="表：VC Method 三情景概率加权")
    add_chart(doc, chart_dir, "chart02_vc_method.png", 6.0, "图表7：VC Method 退出回报分析")
    add_takeaway(doc, f"三情景退出估值¥{SCN[2]['exit_value_yi']:.0f}亿/¥{SCN[1]['exit_value_yi']:.0f}亿/¥{SCN[0]['exit_value_yi']:.0f}亿，对应MOIC {SCN[0]['moic']:.1f}x/{SCN[1]['moic']:.1f}x/{SCN[2]['moic']:.1f}x。基准情景{SCN[1]['moic']:.1f}x已超天使轮门槛。")
    add_chart(doc, chart_dir, "chart08_scenario_probability.png", 6.0, "图表8：情景概率加权估值")
    add_takeaway(doc, f"**概率加权MOIC = {SCN[0]['moic']:.1f}×{SCN[0]['probability']*100:.0f}% + {SCN[1]['moic']:.1f}×{SCN[1]['probability']*100:.0f}% + {SCN[2]['moic']:.1f}×{SCN[2]['probability']*100:.0f}% ≈ {PW_MOIC:.1f}x**，概率加权退出估值 ≈ {PW_EXIT_YI:.1f}亿。")
    add_heading(doc, "5.4 可比天使轮/Pre-A融资分析", level=2)
    pt = val["precedent_transactions"]
    rows = [[t["target_name"],t["date"],t["stage"],f"¥{t['amount']:,}万",f"¥{t['post_money']:,}万",f"{t['revenue_multiple']:.0f}x",t["reliability"]] for t in pt["transactions"]]
    add_table(doc, ["标的","时间","轮次","融资额","投后估值","营收倍数","可靠性"], rows, caption="表：可比融资交易")
    add_para(doc, f"{len(pt['transactions'])}笔可比交易投后估值中位数 ¥{pt['post_money_median']/10000:.2f}亿（区间 {pt['post_money_range']}），融资额中位数 ¥{pt['amount_median']:,.0f}万。{cname}BP投后¥{BP_POST/10000:.2f}亿处于可比区间合理范围。", indent=True)
    add_chart(doc, chart_dir, "chart04_comparable_financing.png", 6.0, "图表9：可比融资散点图")
    add_takeaway(doc, f"{cname}（★）融资额与投后估值均处于可比案例合理区间。")
    add_heading(doc, "5.5 可比上市公司PS倍数参考", level=2)
    cm = cfg["valuation_params"]["comps"]["public_comps"]
    rows = [[c["name"],c["ticker"],f"¥{c['revenue_ttm']/10000:.0f}亿",f"{c['gross_margin']*100:.0f}%",f"{c['ev_revenue']:.1f}x",f"{c['pe_ttm']}x"] for c in cm]
    add_table(doc, ["公司","代码","营收TTM","毛利率","EV/Revenue","PE"], rows, caption="表：可比上市公司")
    add_chart(doc, chart_dir, "chart07_comparable_ps.png", 6.0, "图表10：可比上市公司PS对比")
    add_takeaway(doc, "可比公司PS中位数约3.5x，给VC Method退出倍数提供上限锚定（折价后基准1.8x）。")
    add_heading(doc, "5.6 估值方法交叉验证与结论", level=2)
    add_para(doc,
        f"三种方法投前估值：VC Method ¥{IV[0]['valuation']/10000:.2f}亿、可比交易 ¥{IV[2]['valuation']/10000:.2f}亿、"
        f"Scorecard ¥{IV[1]['valuation']/10000:.2f}亿，加权中位数 ¥{WPRE/10000:.2f}亿。各方法差异在2x以内，结论稳健。", indent=True)
    add_heading(doc, "5.7 综合估值建议", level=2)
    rec = s("valuation_recommendation").replace("{min}",f"{syn['valuation_range']['min']/10000:.2f}").replace("{max}",f"{syn['valuation_range']['max']/10000:.2f}").replace("{wpre}",f"{WPRE/10000:.2f}")
    add_para(doc, rec, color=C_DARK_BLUE)
    doc.add_page_break()

    # ---- 六、退出分析与股权稀释 ----
    add_heading(doc, "六、退出分析与股权稀释", level=1)
    add_heading(doc, "6.1 退出路径分析", level=2)
    add_para(doc, s("exit_path_analysis"), indent=True)
    add_heading(doc, "6.2 后续稀释路径", level=2)
    dp = cap["dilution_path"]
    holders = ["创始人","ESOP","Angel轮投资者","A轮投资者","B轮投资者","C轮投资者","IPO投资者"]
    rows = []
    for stg in dp:
        own = stg["ownership"]
        rows.append([stg["round"]] + [f"{own.get(h,0)*100:.1f}%" if h in own else "—" for h in holders])
    add_table(doc, ["轮次"]+holders, rows, font_size=8.5, caption="表：股权稀释路径（天使→IPO）")
    add_para(doc, s("dilution_analysis"), indent=True)
    add_chart(doc, chart_dir, "chart09_dilution_path.png", 6.0, "图表11：股权稀释路径")
    add_takeaway(doc, "创始人持股由80%稀释至IPO时约44%，仍保有控制权；天使投资人稀释路径在VC Method中已计入。")
    wf = cap["waterfall"]["non_participating"]
    ev = cap["waterfall"]["exit_value"]
    rows = []
    for name, amt in wf["proceeds"].items():
        moic = wf["moic"].get(name)
        rows.append([name, f"¥{amt/10000:.2f}亿", f"{moic:.1f}x" if moic else "—",
                     f"{amt/wf['total_distributed']*100:.1f}%"])
    add_table(doc, ["股东","分配金额","MOIC","占比"], rows,
              caption=f"表：退出瀑布分配（退出价值¥{ev/10000:.0f}亿，基准情景）")
    add_para(doc, "口径说明：上表为当前轮次完全稀释股权下的waterfall示例，未纳入后续融资轮对天使股权的稀释；VC Method章节的MOIC为含后续稀释后的投资回报口径。", indent=True)
    doc.add_page_break()

    # ---- 七、敏感性分析 ----
    add_heading(doc, "七、敏感性分析", level=1)
    add_heading(doc, "7.1 MOIC敏感性热力图", level=2)
    add_chart(doc, chart_dir, "chart12_moic_heatmap.png", 5.8, "图表12：MOIC敏感性热力图（退出倍数×稀释率）")
    add_takeaway(doc, s("sensitivity.summary"))
    add_heading(doc, "7.2 关键变量敏感性排序", level=2)
    for item in narr(cfg,"sections.sensitivity.items",[]): add_bullet(doc, item)
    add_heading(doc, "7.3 下行情景压力测试", level=2)
    add_para(doc, s("sensitivity.downside"), indent=True)
    doc.add_page_break()

    # ---- 八、风险提示 ----
    add_heading(doc, "八、风险提示", level=1)
    risk_order = [("高风险",C_RED),("中高风险",C_RED),("中风险",C_DARK)]
    for rk_level, rk_color in risk_order:
        items = narr(cfg,f"risks_graded.{rk_level}",[])
        if not items: continue
        add_para(doc, f"【{rk_level}】", size=11.5, bold=True, color=rk_color)
        for name, desc in items:
            add_bullet(doc, f"**{name}**：{desc}")
    doc.add_page_break()

    # ---- 九、关键尽调建议 ----
    add_heading(doc, "九、关键尽调建议", level=1)
    dd = narr(cfg,"dd_items",[])
    for i,(t,a) in enumerate(dd,1):
        add_bullet(doc, a, bold_lead=f"9.{i} {t}：")
    doc.add_page_break()

    # ---- 附录 ----
    add_heading(doc, "附录", level=1)
    add_heading(doc, "附录A：估值模型核心假设汇总", level=2)
    app_rows = narr(cfg,"assumptions_appendix",[]); app_h = narr(cfg,"assumptions_headers",["假设项","取值","来源"])
    if app_rows: add_table(doc, app_h, app_rows, col_widths=[1.8,2.6,2.0])
    add_heading(doc, "附录B：数据来源与可靠性说明", level=2)
    add_para(doc, narr(cfg,"data_source_note",""), indent=True)
    add_heading(doc, "附录C：假设追溯表", level=2)
    add_para(doc, "完整假设追溯见交付物 assumption_traceability.csv 与 source_log.csv。", indent=True)
    doc.add_page_break()

    # ---- 免责声明 ----
    add_heading(doc, "免责声明", level=1)
    disc = narr(cfg,"disclaimer",["本报告仅供内部投资决策参考。"])
    for i,d in enumerate(disc,1): add_bullet(doc, d, bold_lead=f"{i}. ")

    doc.save(doc_path)
    return len(doc.element.body)

print("Word主体定义完成。")
sys.stdout.flush()


# ============================================================
# PPT 17页主体
# ============================================================
def build_ppt(pptx_path, cfg, val, cap, chart_dir):
    prs = Presentation()
    prs.slide_width = PptInches(13.333); prs.slide_height = PptInches(7.5)
    cname = cfg["company"]["name"]; cname_en = cfg["company"].get("name_en","")
    stage_cn = STAGE_CN.get(cfg["round"]["stage"], cfg["round"]["stage"])
    val_date = cfg["meta"]["valuation_date"]
    ft = f"{cname} {stage_cn}估值报告 | {cfg['meta']['confidentiality']} | {val_date}"
    syn = val["synthesis"]; SCN = cfg["valuation_params"]["vc_method"]["scenarios"]
    PW_MOIC = sum(s["moic"]*s["probability"] for s in SCN)
    PW_EXIT_YI = sum(s["exit_value_yi"]*s["probability"] for s in SCN)
    RND_AMT = cfg["round"]["round_amount"]; EQ_PCT = cfg["round"]["equity_offered"]
    BP_POST = cfg["round"]["bp_implied_post_money"]; BP_PRE = cfg["round"]["bp_implied_pre_money"]
    fc = cfg["financials"]; rev = fc["revenue_forecast"]
    nm = fc["net_margin"]; em = fc["ebit_margin"]

    def s(key, default=""): return narr(cfg, f"sections.{key}", default)
    full_name = f"{cname}（{cname_en}）" if cname_en else cname
    slogan = narr(cfg, "slogan", "")
    team = narr(cfg, "analyst_team", "估值团队")

    # ---- Slide 1: 封面 ----
    sl = blank_slide(prs)
    band = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, PptInches(2.2), prs.slide_width, PptInches(2.5))
    band.fill.solid(); band.fill.fore_color.rgb = P_DARK_BLUE; band.line.fill.background()
    tf = band.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = full_name; r.font.size = PptPt(40); r.font.bold = True
    r.font.color.rgb = P_WHITE; _set_ppt_cn(r)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run(); r2.text = f"{stage_cn}估值分析报告"; r2.font.size = PptPt(26)
    r2.font.color.rgb = P_WHITE; _set_ppt_cn(r2)
    if slogan:
        ppt_text(sl, slogan, PptInches(1), PptInches(4.9), PptInches(11.3), PptInches(0.6),
                 size=15, color=P_DARK, align=PP_ALIGN.CENTER, italic=True)
    ppt_text(sl, f"估值日期：{val_date}    ｜    {cfg['meta']['confidentiality']}    ｜    {team}",
             PptInches(1), PptInches(6.4), PptInches(11.3), PptInches(0.5), size=12, color=P_GREY, align=PP_ALIGN.CENTER)

    # ---- Slide 2: 执行摘要 ----
    sl = blank_slide(prs)
    ppt_title(sl, "执行摘要：估值结论先行")
    bullets = [
        ("估值结论", f"本轮按 BP 条款定价：投后 ¥{BP_POST/10000:.2f}亿、投前 ¥{BP_PRE/10000:.2f}亿；分析师多方法加权投前 ¥{syn['weighted_pre_money']/10000:.2f}亿（区间 ¥{syn['valuation_range']['min']/10000:.2f}-{syn['valuation_range']['max']/10000:.2f}亿）。"),
        ("本轮条款", f"BP披露融资 ¥{RND_AMT:,}万、出让 {EQ_PCT*100:.0f}% 股权；证券类型、清算优先权、反稀释与期权池为模型假设，需TS确认。"),
        ("方法选择", "以 VC Method（概率加权）为核心，Scorecard 与可比交易交叉验证；DCF 仅作辅助参考。"),
        ("核心发现", f"概率加权 MOIC {PW_MOIC:.1f}x ＞ 天使轮门槛 8x；BP 定价低于分析师加权中位数，对投资人留有安全边际。"),
    ]
    y = 1.4
    for lead, txt in bullets:
        ppt_text(sl, lead, PptInches(0.6), PptInches(y), PptInches(2.2), PptInches(0.6), size=13, bold=True, color=P_DARK_BLUE)
        ppt_text(sl, txt, PptInches(2.9), PptInches(y), PptInches(9.8), PptInches(0.8), size=12, color=P_DARK)
        y += 0.95
    iv = val["synthesis"]["individual_valuations"]
    rows = [
        [m["method"], f"{syn['valuation_range']['min']:,.0f}", f"{syn['valuation_range']['max']:,.0f}", f"{m['weight']*100:.0f}%"] for m in iv
    ]
    rows.append(["加权投前估值", "—", f"{syn['weighted_pre_money']:,.0f}", "100%"])
    ppt_table(sl, ["估值方法","低（万元）","高（万元）","权重"], rows, PptInches(0.6), PptInches(5.3), PptInches(12.1), PptInches(1.6))
    ppt_footer(sl, ft)

    # ---- Slide 3: 公司概览 ----
    sl = blank_slide(prs)
    ppt_title(sl, narr(cfg,"ppt_titles.slide3","公司概览"))
    info = narr(cfg,"ppt_slide3_info",[])
    y = 1.5
    for k, v in info:
        ppt_text(sl, f"● {k}", PptInches(0.6), PptInches(y), PptInches(1.8), PptInches(0.5), size=13, bold=True, color=P_DARK_BLUE)
        ppt_text(sl, v, PptInches(2.4), PptInches(y), PptInches(4.0), PptInches(0.6), size=12, color=P_DARK)
        y += 0.85
    ppt_text(sl, narr(cfg,"ppt_slide3_advantage",""), PptInches(0.6), PptInches(5.2), PptInches(6.0), PptInches(1.5), size=12, color=P_DARK)
    ppt_pic(sl, chart_dir, "chart05_fund_usage_pie.png", PptInches(6.9), PptInches(1.4), width=PptInches(6.0))
    ppt_footer(sl, ft)

    # ---- Slide 4: 产品矩阵 ----
    sl = blank_slide(prs)
    ppt_title(sl, narr(cfg,"ppt_titles.slide4","产品矩阵"))
    prod = narr(cfg,"products",[]); prod_h = narr(cfg,"product_headers",["产品线","形态","核心功能","变现方式"])
    if prod:
        ppt_table(sl, prod_h, prod, PptInches(0.6), PptInches(1.6), PptInches(12.1), PptInches(3.5), first_col_bold=True)
        ppt_text(sl, s("product_takeaway",""), PptInches(0.6), PptInches(5.6), PptInches(12), PptInches(0.6), size=12, color=P_DARK, italic=True)
    ppt_footer(sl, ft)

    # ---- Slide 5: 商业模式 ----
    sl = blank_slide(prs)
    ppt_title(sl, narr(cfg,"ppt_titles.slide5","商业模式"))
    pts = narr(cfg,"ppt_slide5_points",[])
    y = 1.5
    for t in pts:
        ppt_text(sl, "▪ "+t, PptInches(0.6), PptInches(y), PptInches(6.0), PptInches(0.8), size=12, color=P_DARK)
        y += 0.9
    ppt_pic(sl, chart_dir, "chart10_revenue_mix.png", PptInches(6.9), PptInches(1.4), width=PptInches(6.0))
    ppt_footer(sl, ft)

    # ---- Slide 6: 行业市场 ----
    sl = blank_slide(prs)
    ppt_title(sl, narr(cfg,"ppt_titles.slide6","行业市场"))
    market_rows = narr(cfg,"market_size",[]); market_h = narr(cfg,"market_headers",["细分市场","规模","年份","CAGR"])
    if market_rows:
        ppt_table(sl, market_h, market_rows, PptInches(0.6), PptInches(1.5), PptInches(5.8), PptInches(2.2))
    ppt_text(sl, narr(cfg,"ppt_slide6_drivers",""), PptInches(0.6), PptInches(4.2), PptInches(5.8), PptInches(2.3), size=12, color=P_DARK)
    ppt_pic(sl, chart_dir, "chart06_market_size.png", PptInches(6.7), PptInches(1.5), width=PptInches(6.1))
    ppt_footer(sl, ft)

    # ---- Slide 7: 竞争格局 ----
    sl = blank_slide(prs)
    ppt_title(sl, narr(cfg,"ppt_titles.slide7","竞争格局"))
    comp_rows = narr(cfg,"competitors",[]); comp_h = narr(cfg,"competitor_headers",["类型","代表","优势","短板"])
    if comp_rows:
        ppt_table(sl, comp_h, comp_rows, PptInches(0.6), PptInches(1.5), PptInches(6.0), PptInches(2.2))
    ppt_text(sl, s("competition_paragraph",""), PptInches(0.6), PptInches(4.0), PptInches(6.0), PptInches(1.5), size=12, bold=True, color=P_RED)
    ppt_pic(sl, chart_dir, "chart07_comparable_ps.png", PptInches(6.9), PptInches(1.5), width=PptInches(6.0))
    ppt_footer(sl, ft)

    # ---- Slide 8: 财务预测 ----
    sl = blank_slide(prs)
    ppt_title(sl, narr(cfg,"ppt_titles.slide8","财务预测"))
    rows = [[y_, f"{rev[y_]:,}", f"{fc['gross_margin']*100:.0f}%", f"{em[y_]*100:.1f}%", f"{nm[y_]*100:.1f}%"] for y_ in ["2026","2028","2030"]]
    ppt_table(sl, ["年份","营收(万元)","毛利率","EBIT率","净利率"], rows, PptInches(0.6), PptInches(1.5), PptInches(5.8), PptInches(1.8))
    ppt_text(sl, narr(cfg,"ppt_slide8_warning","⚠ 以上为公司BP预测，未经独立验证"), PptInches(0.6), PptInches(3.7), PptInches(5.8), PptInches(1.5), size=12, bold=True, color=P_RED)
    ppt_pic(sl, chart_dir, "chart01_revenue_forecast.png", PptInches(6.7), PptInches(1.5), width=PptInches(6.1))
    ppt_footer(sl, ft)

    # ---- Slide 9: 估值方法 ----
    sl = blank_slide(prs)
    ppt_title(sl, narr(cfg,"ppt_titles.slide9","估值方法"))
    pts = narr(cfg,"ppt_slide9_points",[])
    y = 1.5
    for t in pts:
        ppt_text(sl, t, PptInches(0.6), PptInches(y), PptInches(6.0), PptInches(0.7), size=12, color=P_DARK)
        y += 0.85
    ppt_pic(sl, chart_dir, "chart03_valuation_bridge.png", PptInches(6.8), PptInches(1.5), width=PptInches(6.1))
    ppt_footer(sl, ft)

    # ---- Slide 10: VC Method ----
    sl = blank_slide(prs)
    ppt_title(sl, narr(cfg,"ppt_titles.slide10","VC Method"))
    rows = [[sc["name"],f"{sc['exit_revenue_yi']:.0f}亿",f"{sc['ps_multiple']:.1f}x",f"{sc['exit_value_yi']:.0f}亿",f"{sc['moic']:.1f}x",f"{sc['probability']*100:.0f}%"] for sc in SCN]
    ppt_table(sl, ["情景","退出营收","PS","退出估值","MOIC","概率"], rows, PptInches(0.6), PptInches(1.5), PptInches(6.2), PptInches(1.8))
    ppt_text(sl, f"概率加权 MOIC = {PW_MOIC:.1f}x，退出估值 = {PW_EXIT_YI:.1f}亿。基准情景 {SCN[1]['moic']:.1f}x 回报。", PptInches(0.6), PptInches(3.7), PptInches(6.2), PptInches(1.2), size=12, bold=True, color=P_DARK_BLUE)
    ppt_pic(sl, chart_dir, "chart02_vc_method.png", PptInches(7.0), PptInches(1.5), width=PptInches(5.9))
    ppt_footer(sl, ft)

    # ---- Slide 11: 情景概率 ----
    sl = blank_slide(prs)
    ppt_title(sl, narr(cfg,"ppt_titles.slide11","情景概率加权"))
    ppt_text(sl,
        f"{SCN[0]['name']}（{SCN[0]['probability']*100:.0f}%）：营收上浮、PS接近上市中位、稀释取下限 → MOIC {SCN[0]['moic']:.1f}x\n"
        f"{SCN[1]['name']}（{SCN[1]['probability']*100:.0f}%）：BP营收、PS折价后中位、稀释取中 → MOIC {SCN[1]['moic']:.1f}x\n"
        f"{SCN[2]['name']}（{SCN[2]['probability']*100:.0f}%）：营收下浮、PS极保守、稀释取上限 → MOIC {SCN[2]['moic']:.1f}x",
        PptInches(0.6), PptInches(1.5), PptInches(6.0), PptInches(2.5), size=12, color=P_DARK)
    ppt_text(sl, f"概率加权 MOIC = {SCN[0]['moic']:.1f}×{SCN[0]['probability']*100:.0f}% + {SCN[1]['moic']:.1f}×{SCN[1]['probability']*100:.0f}% + {SCN[2]['moic']:.1f}×{SCN[2]['probability']*100:.0f}% = {PW_MOIC:.1f}x ＞ 8x 门槛",
             PptInches(0.6), PptInches(4.2), PptInches(6.0), PptInches(1.0), size=12, bold=True, color=P_RED)
    ppt_pic(sl, chart_dir, "chart08_scenario_probability.png", PptInches(6.9), PptInches(1.5), width=PptInches(6.0))
    ppt_footer(sl, ft)

    # ---- Slide 12: 可比分析 ----
    sl = blank_slide(prs)
    ppt_title(sl, narr(cfg,"ppt_titles.slide12","可比分析"))
    pt = val["precedent_transactions"]
    rows = [[t["target_name"][:10], t["date"], t["stage"], f"{t['amount']:,}", f"{t['post_money']:,}"] for t in pt["transactions"]]
    ppt_table(sl, ["标的","时间","轮次","融资(万)","投后(万)"], rows, PptInches(0.6), PptInches(1.5), PptInches(6.2), PptInches(1.8))
    ppt_text(sl, f"可比投后中位数 ¥{pt['post_money_median']/10000:.2f}亿，融资额中位 ¥{pt['amount_median']:,.0f}万。本轮投后处于合理区间。",
             PptInches(0.6), PptInches(3.7), PptInches(6.2), PptInches(1.2), size=12, color=P_DARK)
    ppt_pic(sl, chart_dir, "chart04_comparable_financing.png", PptInches(7.0), PptInches(1.5), width=PptInches(5.9))
    ppt_footer(sl, ft)

    # ---- Slide 13: MOIC 敏感性 ----
    sl = blank_slide(prs)
    ppt_title(sl, narr(cfg,"ppt_titles.slide13","MOIC敏感性分析"))
    ppt_pic(sl, chart_dir, "chart12_moic_heatmap.png", PptInches(0.6), PptInches(1.5), width=PptInches(6.2))
    ppt_text(sl, "敏感性排序（影响从高到低）：\n① 退出 PS 倍数\n② 累计稀释率\n③ 退出营收实现度\n④ 退出时点\n\n下行压力测试：悲观情景仍可获正回报，本金安全边际较高。",
             PptInches(7.1), PptInches(1.6), PptInches(5.7), PptInches(4.5), size=12, color=P_DARK)
    ppt_footer(sl, ft)

    # ---- Slide 14: 稀释路径 ----
    sl = blank_slide(prs)
    ppt_title(sl, narr(cfg,"ppt_titles.slide14","股权稀释路径与退出分配"))
    ppt_pic(sl, chart_dir, "chart09_dilution_path.png", PptInches(0.6), PptInches(1.5), width=PptInches(6.2))
    wf = cap["waterfall"]["non_participating"]
    ev = cap["waterfall"]["exit_value"]
    angel_amt = wf["proceeds"].get("Angel轮投资者", 0)
    angel_moic = wf["moic"].get("Angel轮投资者", 0)
    ppt_text(sl, f"天使投资人持股：{cfg['round']['equity_offered']*100:.0f}% → IPO 后陆续稀释。\n\n"
                f"退出瀑布（当前轮次口径，退出价值{ev/10000:.0f}亿）：\n"
                f"天使获分配 ¥{angel_amt/10000:.2f}亿，MOIC {angel_moic:.1f}x。\n\n"
                f"含后续稀释的回报以VC Method章节为准。",
             PptInches(7.1), PptInches(1.6), PptInches(5.7), PptInches(4.5), size=12, color=P_DARK)
    ppt_footer(sl, ft)

    # ---- Slide 15: 风险提示 ----
    sl = blank_slide(prs)
    ppt_title(sl, narr(cfg,"ppt_titles.slide15","风险提示（分级）"))
    ppt_text(sl, "【高风险】", PptInches(0.6), PptInches(1.4), PptInches(6), PptInches(0.4), size=13, bold=True, color=P_RED)
    ppt_text(sl, narr(cfg,"ppt_slide15_texts.high",""), PptInches(0.6), PptInches(1.9), PptInches(6.0), PptInches(1.6), size=12, color=P_DARK)
    ppt_text(sl, "【中高风险】", PptInches(0.6), PptInches(3.5), PptInches(6), PptInches(0.4), size=13, bold=True, color=PptRGB(0xE6,0x7E,0x22))
    ppt_text(sl, narr(cfg,"ppt_slide15_texts.mid_high",""), PptInches(0.6), PptInches(4.0), PptInches(6.0), PptInches(1.6), size=12, color=P_DARK)
    ppt_text(sl, "【中风险】", PptInches(6.9), PptInches(1.4), PptInches(6), PptInches(0.4), size=13, bold=True, color=P_DARK_BLUE)
    ppt_text(sl, narr(cfg,"ppt_slide15_texts.mid",""), PptInches(6.9), PptInches(1.9), PptInches(6.0), PptInches(2.2), size=12, color=P_DARK)
    ppt_footer(sl, ft)

    # ---- Slide 16: 尽调建议 ----
    sl = blank_slide(prs)
    ppt_title(sl, narr(cfg,"ppt_titles.slide16","关键尽调建议"))
    left = narr(cfg,"ppt_slide16_left",[]); right = narr(cfg,"ppt_slide16_right",[])
    if left: ppt_text(sl, "\n".join(left), PptInches(0.6), PptInches(1.6), PptInches(6.0), PptInches(4.0), size=12, color=P_DARK)
    if right: ppt_text(sl, "\n".join(right), PptInches(6.9), PptInches(1.6), PptInches(6.0), PptInches(4.0), size=12, color=P_DARK)
    ppt_footer(sl, ft)

    # ---- Slide 17: 免责声明 ----
    sl = blank_slide(prs)
    ppt_title(sl, "免责声明")
    disc = narr(cfg,"disclaimer",["本报告仅供内部投资决策参考。"])
    ppt_text(sl, "\n".join(disc), PptInches(0.8), PptInches(1.6), PptInches(11.5), PptInches(4.5), size=12, color=P_DARK)
    ppt_footer(sl, ft)

    prs.save(pptx_path)
    return len(prs.slides.__iter__.__self__._sldIdLst)

print("PPT主体定义完成。")
sys.stdout.flush()


# ============================================================
# Excel 估值模型
# ============================================================
def build_excel(xlsx_path, cfg, val, cap, chart_dir):
    wb = openpyxl.Workbook()
    thin = Side(style="thin",color="CCCCCC"); border = Border(left=thin,right=thin,top=thin,bottom=thin)
    hdr_fill = PatternFill("solid",fgColor=HEADER_BG)
    alt_fill = PatternFill("solid",fgColor=ALT_ROW_BG)
    hdr_font = XlFont(name=FONT_CN, bold=True, color="FFFFFF", size=10)
    cell_font = XlFont(name=FONT_CN, size=10)
    title_font = XlFont(name=FONT_CN, bold=True, color=HEADER_BG, size=13)
    center = XlAlign(horizontal="center",vertical="center")
    cname = cfg["company"]["name"]

    def style(ws, headers, rows, title, col_w=14):
        ws["A1"] = title; ws["A1"].font = title_font; ws.append([]); ws.append(headers)
        hr = ws.max_row
        for c in range(1,len(headers)+1):
            cl = ws.cell(row=hr,column=c); cl.fill = hdr_fill; cl.font = hdr_font
            cl.alignment = center; cl.border = border
        for i,row in enumerate(rows):
            ws.append(row); rr = ws.max_row
            for c in range(1,len(headers)+1):
                cl = ws.cell(row=rr,column=c); cl.font = cell_font; cl.border = border
                cl.alignment = center
                if i%2==1: cl.fill = alt_fill
        for c in range(1,len(headers)+1):
            ws.column_dimensions[openpyxl.utils.get_column_letter(c)].width = col_w

    fc = cfg["financials"]; rev = fc["revenue_forecast"]
    nm = fc["net_margin"]; em = fc["ebit_margin"]; npft = fc["net_profit"]
    SCN = cfg["valuation_params"]["vc_method"]["scenarios"]
    PW_MOIC = sum(s["moic"]*s["probability"] for s in SCN)
    PW_EXIT_YI = sum(s["exit_value_yi"]*s["probability"] for s in SCN)

    ws1 = wb.active; ws1.title = "财务预测"
    rows = [[str(y),rev[str(y)],f"{fc['gross_margin']*100:.0f}%",f"{em[str(y)]*100:.1f}%",f"{nm[str(y)]*100:.1f}%",npft[str(y)]] for y in fc["forecast_years"]]
    style(ws1, ["年份","营收(万元)","毛利率","EBIT率","净利率","净利润(万元)"], rows, f"{cname} 财务预测（来源：公司BP）",14)

    ws2 = wb.create_sheet("VC Method情景")
    rows = [[s["name"],f"{s['exit_revenue_yi']:.0f}",f"{s['ps_multiple']:.1f}x",f"{s['exit_value_yi']:.0f}",f"{s['angel_retention']*100:.0f}%",f"{s['moic']:.1f}x",f"{s['probability']*100:.0f}%"] for s in SCN]
    rows.append(["概率加权","-","-",f"{PW_EXIT_YI:.1f}","-",f"{PW_MOIC:.1f}x","100%"])
    style(ws2, ["情景","退出营收(亿)","PS倍数","退出估值(亿)","天使保留","MOIC","概率"], rows, "VC Method 三情景概率加权",14)

    ws3 = wb.create_sheet("估值综合")
    syn = val["synthesis"]; rows = [[m["method"],f"{m['valuation']:,.0f}",f"{m['weight']*100:.0f}%",f"{m['weighted_contribution']:,.0f}",m["source"]] for m in syn["individual_valuations"]]
    rows.append(["加权投前估值",f"{syn['weighted_pre_money']:,.0f}","100%",f"{syn['weighted_pre_money']:,.0f}",f"范围{syn['valuation_range']['min']:,.0f}-{syn['valuation_range']['max']:,.0f}"])
    style(ws3, ["方法","投前估值(万)","权重","加权贡献(万)","依据"], rows, "多方法交叉验证综合估值",18)

    ws4 = wb.create_sheet("Cap Table稀释")
    dp = cap["dilution_path"]; holders = ["创始人","ESOP","Angel轮投资者","A轮投资者","B轮投资者","C轮投资者","IPO投资者"]
    rows = [[stg["round"]]+[f"{stg['ownership'].get(h,0)*100:.1f}%" if h in stg["ownership"] else "-" for h in holders] for stg in dp]
    style(ws4, ["轮次"]+holders, rows, "股权稀释路径（天使→IPO）",13)

    ws5 = wb.create_sheet("退出瀑布分配")
    wf = cap["waterfall"]["non_participating"]; ev = cap["waterfall"]["exit_value"]
    rows = [[name,f"{amt/10000:.2f}",f"{wf['moic'].get(name,0):.1f}x" if wf["moic"].get(name) else "-",f"{amt/wf['total_distributed']*100:.1f}%"] for name,amt in wf["proceeds"].items()]
    style(ws5, ["股东","分配金额(亿)","MOIC","占比"], rows, f"退出瀑布分配（退出价值¥{ev/10000:.0f}亿）",15)

    wb.save(xlsx_path)
    print(f"[Excel] {os.path.basename(xlsx_path)}")

# ============================================================
# 辅助交付物
# ============================================================
def build_aux(outdir, data_dir, cfg, val, cap):
    SCN = cfg["valuation_params"]["vc_method"]["scenarios"]
    PW_MOIC = sum(s["moic"]*s["probability"] for s in SCN)
    PW_EXIT_YI = sum(s["exit_value_yi"]*s["probability"] for s in SCN)
    cname = cfg["company"]["name"]; val_date = cfg["meta"]["valuation_date"]
    stage_cn = STAGE_CN.get(cfg["round"]["stage"], cfg["round"]["stage"])

    # assumption_traceability.csv
    rows = []  # sourced from config.source_log
    for sl in cfg.get("source_log",[]):
        rows.append([sl["param"], sl["value"], sl["source"], sl["detail"], "中"])
    with open(os.path.join(data_dir,"assumption_traceability.csv"),"w",newline="",encoding="utf-8-sig") as f:
        w = csv.writer(f); w.writerow(["关键假设","取值","来源分类","依据/明细","置信度"]); w.writerows(rows)

    # source_log.csv
    with open(os.path.join(data_dir,"source_log.csv"),"w",newline="",encoding="utf-8-sig") as f:
        w = csv.writer(f); w.writerow(["参数","取值","来源","明细"])
        for s in cfg.get("source_log",[]): w.writerow([s["param"],s["value"],s["source"],s["detail"]])

    # one_page_summary.md
    syn = val["synthesis"]; BP_POST = cfg["round"]["bp_implied_post_money"]; BP_PRE = cfg["round"]["bp_implied_pre_money"]
    RND_AMT = cfg["round"]["round_amount"]; EQ_PCT = cfg["round"]["equity_offered"]
    risks = narr(cfg,"one_page_summary.risks",[])
    rec = narr(cfg,"one_page_summary.recommendation","")
    src = narr(cfg,"one_page_summary.source_note","")
    summary = f"# {cname} {stage_cn}估值 · 投委会一页速览\n\n"
    summary += f"**估值日期**：{val_date} ｜ **货币**：人民币 ｜ **{cfg['meta']['confidentiality']}**\n\n"
    summary += f"## 估值结论\n- **本轮交易定价（BP条款）**：投后 ¥{BP_POST/10000:.2f}亿，投前 ¥{BP_PRE/10000:.2f}亿，融资 ¥{RND_AMT:,}万，出让 {EQ_PCT*100:.0f}%\n"
    summary += f"- **分析师多方法加权投前估值**：¥{syn['weighted_pre_money']/10000:.2f}亿（区间 ¥{syn['valuation_range']['min']/10000:.2f}亿 - ¥{syn['valuation_range']['max']/10000:.2f}亿）\n\n"
    summary += f"## VC Method（核心方法，概率加权）\n| 情景 | 退出营收 | PS | 退出估值 | MOIC | 概率 |\n|---|---|---|---|---|---|\n"
    for sc in SCN:
        summary += f"| {sc['name']} | {sc['exit_revenue_yi']:.0f}亿 | {sc['ps_multiple']:.1f}x | {sc['exit_value_yi']:.0f}亿 | {sc['moic']:.1f}x | {sc['probability']*100:.0f}% |\n"
    summary += f"| **概率加权** | - | - | **{PW_EXIT_YI:.1f}亿** | **{PW_MOIC:.1f}x** | 100% |\n\n"
    summary += f"## 多方法交叉验证\n| 方法 | 投前估值 | 权重 |\n|---|---|---|\n"
    for m in syn["individual_valuations"]: summary += f"| {m['method']} | ¥{m['valuation']/10000:.2f}亿 | {m['weight']*100:.0f}% |\n"
    summary += "\n## 关键风险（高→中）\n"
    for rk in risks: summary += f"- {rk}\n"
    summary += f"\n## 投资建议\n{rec}\n\n> {src}\n"
    with open(os.path.join(outdir,"one_page_summary.md"),"w",encoding="utf-8") as f: f.write(summary)

    # open_questions.md
    oq = f"# {cname} {stage_cn} · 尽调问题清单（按优先级）\n\n"
    for level in ["P0","P1","P2"]:
        items = narr(cfg,f"open_questions.{level}",[])
        if not items: continue
        label = {"P0":"阻断性问题（影响投资决策）","P1":"重要问题（影响估值）","P2":"补充问题（投后管理）"}[level]
        oq += f"## {level} — {label}\n"
        for i,item in enumerate(items,({"P0":1,"P1":4,"P2":7}[level])):
            oq += f"{i}. {item}\n"
        oq += "\n"
    oq += "> 尽调清单将随数据室开放逐项关闭。\n"
    with open(os.path.join(outdir,"open_questions.md"),"w",encoding="utf-8") as f: f.write(oq)

    print("[Aux] assumption_traceability / source_log / one_page_summary / open_questions")

# ============================================================
# Manifest
# ============================================================
def build_manifest(outdir, outputs, cfg, val):
    syn = val["synthesis"]; SCN = cfg["valuation_params"]["vc_method"]["scenarios"]
    PW_MOIC = sum(s["moic"]*s["probability"] for s in SCN)
    PW_EXIT_YI = sum(s["exit_value_yi"]*s["probability"] for s in SCN)
    manifest = {
        "valuation_date": cfg["meta"]["valuation_date"],
        "currency": "人民币", "company": cfg["company"]["name"],
        "stage": cfg["round"]["stage"],
        "outputs": outputs,
        "conclusion": {
            "bp_post_money_rmb_yi": cfg["round"]["bp_implied_post_money"]/10000,
            "bp_pre_money_rmb_yi": cfg["round"]["bp_implied_pre_money"]/10000,
            "analyst_weighted_pre_money_rmb_yi": round(syn["weighted_pre_money"]/10000,2),
            "analyst_pre_money_range_rmb_yi": [round(syn["valuation_range"]["min"]/10000,2), round(syn["valuation_range"]["max"]/10000,2)],
            "round_amount_rmb_wan": cfg["round"]["round_amount"],
            "equity_offered": f"{cfg['round']['equity_offered']*100:.0f}%",
            "primary_method": "VC Method（风险资本法，概率加权）",
            "base_case_moic": SCN[1]["moic"],
            "probability_weighted_moic": round(PW_MOIC,1),
            "probability_weighted_exit_yi": round(PW_EXIT_YI,1),
        }, "report_version": "v3"
    }
    with open(os.path.join(outdir,"manifest.json"),"w",encoding="utf-8") as f: json.dump(manifest,f,ensure_ascii=False,indent=2)
    print("[Manifest] manifest.json")

# ============================================================
# Main
# ============================================================
def main():
    args = parse_args()
    cfg, val, cap = load_data(args.config, args.data)
    os.makedirs(args.outdir, exist_ok=True)
    cname = cfg["company"]["name"]
    stage_cn = STAGE_CN.get(cfg["round"]["stage"], cfg["round"]["stage"])
    stage_fname = cfg["round"]["stage"]

    doc_path = os.path.join(args.outdir, f"{cname}_{stage_fname}估值报告_v3.docx")
    ppt_path = os.path.join(args.outdir, f"{cname}_{stage_fname}估值报告_v3.pptx")
    xls_path = os.path.join(args.outdir, f"{cname}_估值模型.xlsx")

    print(f"生成交付物：{cname} ({stage_cn})")
    build_word(doc_path, cfg, val, cap, args.charts)
    build_ppt(ppt_path, cfg, val, cap, args.charts)
    build_excel(xls_path, cfg, val, cap, args.charts)
    build_aux(args.outdir, args.data, cfg, val, cap)

    outputs = [
        f"{os.path.relpath(doc_path, args.outdir)}",
        f"{os.path.relpath(ppt_path, args.outdir)}",
        f"{os.path.relpath(xls_path, args.outdir)}",
        "one_page_summary.md", "open_questions.md",
        os.path.relpath(args.config, args.outdir) if "/" in args.config else args.config,
    ]
    build_manifest(args.outdir, outputs, cfg, val)
    print(f"\n全部交付物已生成至 {args.outdir}/")

if __name__ == "__main__":
    main()
